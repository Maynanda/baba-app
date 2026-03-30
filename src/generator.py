import json
import os
import glob
import subprocess
from pptx import Presentation
from pdf2image import convert_from_path

# Paths
BASE_DIR = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
DATA_DIR = os.path.join(BASE_DIR, "data")
TEMPLATE_PATH = os.path.join(BASE_DIR, "templates", "main_carousel.pptx")
OUTPUT_PPTX_DIR = os.path.join(BASE_DIR, "output", "pptx")
OUTPUT_PDF_DIR = os.path.join(BASE_DIR, "output", "pdf")
OUTPUT_IMG_DIR = os.path.join(BASE_DIR, "output", "images")
SOFFICE_PATH = "/Applications/LibreOffice.app/Contents/MacOS/soffice"

# Ensure output directories exist
os.makedirs(OUTPUT_PPTX_DIR, exist_ok=True)
os.makedirs(OUTPUT_PDF_DIR, exist_ok=True)
os.makedirs(OUTPUT_IMG_DIR, exist_ok=True)

def replace_text_in_shape(shape, placeholder_dict):
    """
    Search and replace text in a PPTX shape based on a dictionary of placeholders.
    Allows formatting to be kept if possible.
    """
    if not shape.has_text_frame:
        return
        
    for paragraph in shape.text_frame.paragraphs:
        para_text = "".join(run.text for run in paragraph.runs)
        
        for key, value in placeholder_dict.items():
            placeholder = f"{{{{{key}}}}}" # e.g. {{TITLE}}
            if placeholder in para_text:
                if paragraph.runs:
                    paragraph.runs[0].text = para_text.replace(placeholder, value)
                    for i in range(1, len(paragraph.runs)):
                        paragraph.runs[i].text = ""

def generate_carousel(json_file_path):
    if not os.path.exists(TEMPLATE_PATH):
        print(f"Error: Template not found at {TEMPLATE_PATH}")
        return

    # Load data
    with open(json_file_path, 'r') as f:
        data = json.load(f)
        
    print(f"Loaded data for post date: {data.get('post_date')}")

    prs = Presentation(TEMPLATE_PATH)
    
    placeholders = {}
    for slide_data in data.get('slides', []):
        for key, val in slide_data.items():
            if key != 'type':
                placeholders[key] = str(val)
                
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                replace_text_in_shape(shape, placeholders)
                
    # 1. Save PPTX
    output_filename = f"carousel_{os.path.basename(json_file_path).replace('.json', '.pptx')}"
    output_path = os.path.join(OUTPUT_PPTX_DIR, output_filename)
    prs.save(output_path)
    print(f"Successfully generated carousel: {output_path}")
    
    # 2. Convert to PDF for LinkedIn
    print("Converting PPTX to PDF (for LinkedIn)...")
    subprocess.run([
        SOFFICE_PATH,
        "--headless",
        "--convert-to", "pdf",
        output_path,
        "--outdir", OUTPUT_PDF_DIR
    ], stdout=subprocess.DEVNULL, stderr=subprocess.DEVNULL)
    
    pdf_filename = output_filename.replace('.pptx', '.pdf')
    pdf_path = os.path.join(OUTPUT_PDF_DIR, pdf_filename)
    print(f"Successfully generated PDF: {pdf_path}")
    
    # 3. Convert PDF to Images for TikTok / Instagram
    print("Converting PDF to Images (for TikTok/IG)...")
    post_date = data.get('post_date', 'unknown_date')
    post_slug = os.path.basename(json_file_path).replace('.json', '')
    post_img_dir = os.path.join(OUTPUT_IMG_DIR, f"{post_date}_{post_slug}")
    os.makedirs(post_img_dir, exist_ok=True)
    
    images = convert_from_path(pdf_path)
    for i, img in enumerate(images):
        img_path = os.path.join(post_img_dir, f"slide_{i+1}.png")
        img.save(img_path, "PNG")
    print(f"Successfully generated {len(images)} images in: {post_img_dir}")

def process_all_files():
    json_files = glob.glob(os.path.join(DATA_DIR, "*.json"))
    for file in json_files:
        print(f"Processing {file}...")
        generate_carousel(file)

if __name__ == "__main__":
    process_all_files()

