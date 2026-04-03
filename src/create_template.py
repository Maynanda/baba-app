import os
from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

BASE_DIR = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))

# Colors matching the strategy
DARK_BG = RGBColor(20, 25, 30) # Very dark slate
TEAL = RGBColor(0, 128, 128)
CORAL = RGBColor(255, 127, 80)
TEXT_WHITE = RGBColor(240, 240, 240)
TEXT_GRAY = RGBColor(170, 170, 170)

def apply_dark_bg(slide):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = DARK_BG

def add_footer(slide, slide_num, total_slides, width_in, height_in):
    # Author handle
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(height_in - 0.8), Inches(3), Inches(0.5))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "@DataScientistExplains"
    p.font.color.rgb = TEXT_GRAY
    p.font.size = Pt(16)
    
    # Progress/Slide num
    txBox2 = slide.shapes.add_textbox(Inches(width_in - 1.5), Inches(height_in - 0.8), Inches(1), Inches(0.5))
    tf2 = txBox2.text_frame
    p2 = tf2.paragraphs[0]
    p2.text = f"{slide_num}/{total_slides}"
    p2.font.color.rgb = TEAL
    p2.font.size = Pt(16)

def create_template(width_in: float, height_in: float, output_path: str):
    output_dir = os.path.dirname(output_path)
    if not os.path.exists(output_dir):
        os.makedirs(output_dir)
        
    prs = Presentation()
    prs.slide_width = Inches(width_in)
    prs.slide_height = Inches(height_in)
    
    # Scale factors based on 10x10 baseline width
    scale = width_in / 10.0
    
    # Slide 1: Hook
    slide1 = prs.slides.add_slide(prs.slide_layouts[6]) # blank layout
    apply_dark_bg(slide1)
    add_footer(slide1, 1, 5, width_in, height_in)
    
    # Accent shape
    shape = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.0 * scale), Inches(height_in * 0.1), Inches(3 * scale), Inches(0.1))
    shape.fill.solid()
    shape.fill.fore_color.rgb = CORAL
    shape.line.color.rgb = CORAL
    
    # Hook Title
    txBox = slide1.shapes.add_textbox(Inches(1.0 * scale), Inches(height_in * 0.2), Inches(width_in - (2 * scale)), Inches(4 * scale))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "{{HOOK_TITLE}}"
    p.font.size = Pt(int(64 * scale))
    p.font.bold = True
    p.font.color.rgb = TEXT_WHITE
    
    # Hook Subtitle
    txBox2 = slide1.shapes.add_textbox(Inches(1.0 * scale), Inches(height_in * 0.6), Inches(width_in - (2 * scale)), Inches(1.5 * scale))
    tf2 = txBox2.text_frame
    tf2.word_wrap = True
    p2 = tf2.paragraphs[0]
    p2.text = "{{HOOK_SUB}}"
    p2.font.size = Pt(int(32 * scale))
    p2.font.color.rgb = TEAL
    
    # Slide 2-4: Body
    for slide_idx in range(2, 5):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        apply_dark_bg(slide)
        add_footer(slide, slide_idx, 5, width_in, height_in)
        
        txBox = slide.shapes.add_textbox(Inches(1.0 * scale), Inches(height_in * 0.1), Inches(4 * scale), Inches(0.5))
        p = txBox.text_frame.paragraphs[0]
        p.text = "HARD-WON LESSONS"
        p.font.size = Pt(int(14 * scale))
        p.font.color.rgb = CORAL
        p.font.bold = True
        
        txBox3 = slide.shapes.add_textbox(Inches(1.0 * scale), Inches(height_in * 0.2), Inches(width_in - (2 * scale)), Inches(1.5 * scale))
        p3 = txBox3.text_frame.paragraphs[0]
        p3.text = f"{{{{BODY_{slide_idx - 1}_TITLE}}}}"
        p3.font.size = Pt(int(44 * scale))
        p3.font.bold = True
        p3.font.color.rgb = TEXT_WHITE
        
        txBox4 = slide.shapes.add_textbox(Inches(1.0 * scale), Inches(height_in * 0.4), Inches(width_in - (2 * scale)), Inches(5 * scale))
        txBox4.text_frame.word_wrap = True
        p4 = txBox4.text_frame.paragraphs[0]
        p4.text = f"{{{{BODY_{slide_idx - 1}_TEXT}}}}"
        p4.font.size = Pt(int(28 * scale))
        p4.font.color.rgb = TEXT_GRAY

    # Slide 5: CTA
    slide5 = prs.slides.add_slide(prs.slide_layouts[6])
    apply_dark_bg(slide5)
    add_footer(slide5, 5, 5, width_in, height_in)
    
    shape = slide5.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.0 * scale), Inches(height_in * 0.1), Inches(3 * scale), Inches(0.1))
    shape.fill.solid()
    shape.fill.fore_color.rgb = TEAL
    shape.line.color.rgb = TEAL
    
    txBox3 = slide5.shapes.add_textbox(Inches(1.0 * scale), Inches(height_in * 0.2), Inches(width_in - (2 * scale)), Inches(2 * scale))
    p3 = txBox3.text_frame.paragraphs[0]
    p3.text = "{{CTA_TITLE}}"
    p3.font.size = Pt(int(54 * scale))
    p3.font.bold = True
    p3.font.color.rgb = TEXT_WHITE
    
    txBox4 = slide5.shapes.add_textbox(Inches(1.0 * scale), Inches(height_in * 0.45), Inches(width_in - (2 * scale)), Inches(3 * scale))
    txBox4.text_frame.word_wrap = True
    p4 = txBox4.text_frame.paragraphs[0]
    p4.text = "{{CTA_TEXT}}"
    p4.font.size = Pt(int(32 * scale))
    p4.font.color.rgb = TEXT_GRAY
    
    prs.save(output_path)
    print(f"Template successfully created at {output_path}")

if __name__ == "__main__":
    # LinkedIn / Instagram Feed (1:1)
    create_template(10.0, 10.0, os.path.join(BASE_DIR, "templates", "carousel_dark_1x1", "main_carousel.pptx"))
    # Instagram Story / TikTok (9:16)
    create_template(5.625, 10.0, os.path.join(BASE_DIR, "templates", "story_dark_9x16", "main_story.pptx"))
