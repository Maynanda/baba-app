"""
src/generator/base.py
Shared visual generation logic (PPTX filling, PDF conversion, image extraction).
"""

import os
import json
import subprocess
from pathlib import Path

from pptx import Presentation
from pdf2image import convert_from_path

BASE_DIR = Path(__file__).resolve().parent.parent.parent
SOFFICE_PATH = "/Applications/LibreOffice.app/Contents/MacOS/soffice"
TEMPLATES_REGISTRY = BASE_DIR / "templates" / "registry.json"


def get_template_path(template_id: str) -> Path:
    """Find the PPTX file for a template id."""
    # Priority 1: templates/{id}/template.pptx
    t_dir = BASE_DIR / "templates" / template_id
    if t_dir.exists():
        for candidate in ["template.pptx", "main_carousel.pptx"]:
            if (t_dir / candidate).exists():
                return t_dir / candidate
    
    # Priority 2: root templates/main_carousel.pptx (fallback)
    root_fallback = BASE_DIR / "templates" / "main_carousel.pptx"
    if root_fallback.exists():
        return root_fallback
        
    raise FileNotFoundError(f"No PPTX template found for '{template_id}'.")


def load_template_metadata(template_id: str) -> dict:
    """Read template metadata. Prefers the local template.json inside the folder."""
    # First: Check templates/{id}/template.json
    t_dir = BASE_DIR / "templates" / template_id
    tj = t_dir / "template.json"
    if tj.exists():
        with open(tj, "r") as f:
            return json.load(f)
            
    # Second: Fallback to central registry.json
    if not TEMPLATES_REGISTRY.exists():
        raise FileNotFoundError(f"Registry not found: {TEMPLATES_REGISTRY}")
        
    with open(TEMPLATES_REGISTRY, "r") as f:
        registry = json.load(f)
        
    for t in registry.get("templates", []):
        if t["id"] == template_id:
            return t
            
    raise ValueError(f"Template '{template_id}' not found.")


def replace_text_in_shape(shape, placeholder_dict):
    """
    Search and replace text in a PPTX shape based on a dictionary of placeholders.
    Allows formatting to be kept if possible.
    """
    if not shape.has_text_frame:
        return
        
    for paragraph in shape.text_frame.paragraphs:
        para_text = "".join(run.text for run in paragraph.runs)
        
        for key, value in placeholder_dict.items():
            placeholder = f"{{{{{key}}}}}" # e.g. {{HOOK_TITLE}}
            if placeholder in para_text:
                if paragraph.runs:
                    paragraph.runs[0].text = para_text.replace(placeholder, str(value))
                    for i in range(1, len(paragraph.runs)):
                        paragraph.runs[i].text = ""


def fill_placeholders(pptx_path: Path, data: dict, output_path: Path):
    """Fill placeholder text on the template slides using the data."""
    if not pptx_path.exists():
        raise FileNotFoundError(f"Template file not found: {pptx_path}")

    prs = Presentation(str(pptx_path))
    placeholders = {}
    
    # If 'slides' is a dict (new format), iterate its items. 
    # If list (old format), iterate slides sequentially.
    slides_data = data.get('slides', {})
    if isinstance(slides_data, dict):
        for key, val in slides_data.items():
            placeholders[key] = str(val)
    elif isinstance(slides_data, list):
        for slide_obj in slides_data:
            if isinstance(slide_obj, dict):
                for key, val in slide_obj.items():
                    if key != 'type':
                        placeholders[key] = str(val)
                
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                replace_text_in_shape(shape, placeholders)
                
    # Ensure parent output directory exists
    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output_path))
    return output_path


def export_pdf(pptx_path: Path, output_dir: Path) -> Path:
    """Convert PPTX to PDF using headless LibreOffice."""
    output_dir.mkdir(parents=True, exist_ok=True)
    
    subprocess.run([
        SOFFICE_PATH,
        "--headless",
        "--convert-to", "pdf",
        str(pptx_path),
        "--outdir", str(output_dir)
    ], stdout=subprocess.DEVNULL, stderr=subprocess.DEVNULL)
    
    pdf_filename = pptx_path.name.replace('.pptx', '.pdf')
    pdf_path = output_dir / pdf_filename
    
    if not pdf_path.exists():
        raise RuntimeError(f"PDF conversion failed for {pptx_path}")
        
    return pdf_path


def export_images(pdf_path: Path, output_dir: Path) -> list[Path]:
    """Convert PDF to PNG slides using pdf2image."""
    output_dir.mkdir(parents=True, exist_ok=True)
    
    images = convert_from_path(str(pdf_path))
    saved_paths = []
    
    for i, img in enumerate(images):
        img_path = output_dir / f"slide_{i+1}.png"
        img.save(str(img_path), "PNG")
        saved_paths.append(img_path)
        
    return saved_paths
